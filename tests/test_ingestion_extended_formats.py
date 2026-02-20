import json

import pytest

from app.config import runtime_settings as settings
from app.ingestion import vector_ingest_service
from app.ingestion.extractors import build_default_registry
from app.ingestion.pipeline import _context, build_options, ingest_paths, ingest_uploaded_files


def _fake_embed(*args, **kwargs):
    return {"embedding": [0.1, 0.2, 0.3, 0.4], "model": "fake", "used_fallback": False}


def _configure_store(tmp_path):
    settings.CONFIG["vector_db_persist_dir"] = str(tmp_path / "chroma")
    settings.CONFIG["vector_db_collection"] = "test_ingestion_extended"
    settings.CONFIG["embedding_dim"] = 4
    settings.CONFIG["vector_db_batch_size"] = 8


def test_registry_supports_special_names():
    registry = build_default_registry()
    docker = registry.resolve("Dockerfile")
    makefile = registry.resolve("Makefile")
    assert docker.name == "textual"
    assert makefile.name == "textual"


def test_ingest_paths_for_must_have_formats(tmp_path, monkeypatch):
    _configure_store(tmp_path)
    monkeypatch.setattr(vector_ingest_service, "embed_text", _fake_embed)

    paths = []

    md = tmp_path / "README.md"
    md.write_text("# Title\n\nSome markdown content.", encoding="utf-8")
    paths.append(str(md))

    yaml_file = tmp_path / "config.yaml"
    yaml_file.write_text("name: demo\nversion: 1", encoding="utf-8")
    paths.append(str(yaml_file))

    jsonl = tmp_path / "records.jsonl"
    jsonl.write_text('{"id":1,"text":"hello"}\n{"id":2,"text":"world"}\n', encoding="utf-8")
    paths.append(str(jsonl))

    csv_file = tmp_path / "table.csv"
    csv_file.write_text("col_a,col_b\nfoo,bar\n", encoding="utf-8")
    paths.append(str(csv_file))

    xml_file = tmp_path / "schema.xml"
    xml_file.write_text("<root><item>alpha</item><item>beta</item></root>", encoding="utf-8")
    paths.append(str(xml_file))

    docker = tmp_path / "Dockerfile"
    docker.write_text("FROM python:3.11\nRUN echo hello", encoding="utf-8")
    paths.append(str(docker))

    options = build_options(max_bytes=1024 * 1024, max_rows=100, max_objects=100)
    summary = ingest_paths(paths, options=options)

    assert summary["total_files"] == len(paths)
    assert summary["failed"] == 0
    assert summary["extracted"] >= 5
    assert summary["total_chunks"] > 0


def test_ingest_size_limit_skips_file(tmp_path, monkeypatch):
    _configure_store(tmp_path)
    monkeypatch.setattr(vector_ingest_service, "embed_text", _fake_embed)

    large = tmp_path / "large.txt"
    large.write_text("x" * 4096, encoding="utf-8")

    summary = ingest_paths([str(large)], options=build_options(max_bytes=128))
    assert summary["total_files"] == 1
    assert summary["extracted"] == 0
    assert summary["files"][0]["status"] in {"skipped", "failed"}


def test_upload_ingestion_uses_same_pipeline(tmp_path, monkeypatch):
    _configure_store(tmp_path)
    monkeypatch.setattr(vector_ingest_service, "embed_text", _fake_embed)

    uploaded = [
        ("notes.md", b"# Notes\n\nhello from upload"),
        ("data.json", json.dumps({"name": "demo", "ok": True}).encode("utf-8")),
    ]
    summary = ingest_uploaded_files(uploaded, options=build_options())

    assert summary["total_files"] == 2
    assert summary["failed"] == 0
    assert summary["extracted"] >= 1


def test_office_happy_path_extractors(tmp_path):
    pytest.importorskip("docx")
    pytest.importorskip("openpyxl")
    pytest.importorskip("pptx")

    from docx import Document
    from openpyxl import Workbook
    from pptx import Presentation
    from PyPDF2 import PdfWriter

    registry = build_default_registry()
    context = _context(build_options(max_pages=10, max_rows=100, max_slides=10, max_sheets=5))

    docx_path = tmp_path / "sample.docx"
    document = Document()
    document.add_paragraph("Docx content fixture")
    document.save(docx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    workbook = Workbook()
    ws = workbook.active
    ws["A1"] = "hello"
    ws["B1"] = "world"
    workbook.save(xlsx_path)

    pptx_path = tmp_path / "sample.pptx"
    presentation = Presentation()
    slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Fixture slide"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Slide body"
    presentation.save(pptx_path)

    pdf_path = tmp_path / "sample.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=300, height=300)
    writer.add_metadata({"/Title": "PDF Fixture Text Extract"})
    with open(pdf_path, "wb") as handle:
        writer.write(handle)

    docx_doc = registry.extract_from_path(str(docx_path), context)
    xlsx_doc = registry.extract_from_path(str(xlsx_path), context)
    pptx_doc = registry.extract_from_path(str(pptx_path), context)
    pdf_doc = registry.extract_from_path(str(pdf_path), context)

    assert docx_doc.units
    assert xlsx_doc.units
    assert pptx_doc.units
    assert pdf_doc.units


def test_parquet_enabled_by_default(tmp_path):
    pa = pytest.importorskip("pyarrow")
    pq = pytest.importorskip("pyarrow.parquet")

    registry = build_default_registry()
    context = _context(build_options(max_rows=10))

    table = pa.table({"name": ["alpha", "beta"], "score": [1, 2]})
    parquet_path = tmp_path / "sample.parquet"
    pq.write_table(table, parquet_path)

    doc = registry.extract_from_path(str(parquet_path), context)
    assert doc.units


def test_legacy_office_graceful_failure(tmp_path, monkeypatch):
    _configure_store(tmp_path)
    monkeypatch.setattr(vector_ingest_service, "embed_text", _fake_embed)

    uploaded = [("legacy.doc", b"not-an-ole-document")]
    summary = ingest_uploaded_files(uploaded, options=build_options(enable_legacy_office=True))
    assert summary["total_files"] == 1
    assert summary["failed"] == 0
    assert summary["files"][0]["status"] in {"ok", "skipped"}
