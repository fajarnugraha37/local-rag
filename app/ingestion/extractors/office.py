from __future__ import annotations

import io
import os
import re
import tempfile
from contextlib import contextmanager
from typing import Generator, List, Optional

from .base import ExtractedDocument, ExtractedUnit, ExtractorContext, MissingDependencyError
from .utils import check_zip_safety, normalize_text, read_bytes


@contextmanager
def _bytes_as_temp_file(raw_bytes: Optional[bytes], *, suffix: str) -> Generator[Optional[str], None, None]:
    if raw_bytes is None:
        yield None
        return
    with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
        tmp.write(raw_bytes)
        tmp_path = tmp.name
    try:
        yield tmp_path
    finally:
        try:
            os.remove(tmp_path)
        except OSError:
            pass


def _load_raw(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> bytes:
    if raw_bytes is not None:
        if len(raw_bytes) > context.max_bytes:
            raise ValueError(f"file exceeds max bytes ({len(raw_bytes)} > {context.max_bytes})")
        return raw_bytes
    return read_bytes(path, max_bytes=context.max_bytes)


def _extract_pdf_from_reader(reader, context: ExtractorContext, raw_pdf: bytes) -> ExtractedDocument:
    warnings: List[str] = []
    units: List[ExtractedUnit] = []

    total_pages = len(reader.pages)
    for idx, page in enumerate(reader.pages, start=1):
        if idx > context.max_pages:
            warnings.append(f"pdf_truncated_at_max_pages={context.max_pages}")
            break
        text = normalize_text(page.extract_text() or "")
        if text:
            units.append(ExtractedUnit(text=text, metadata={"page_number": idx}))

    if total_pages > context.max_pages:
        warnings.append(f"pdf_total_pages={total_pages}")

    if not units:
        # Best-effort fallback for PDFs where page-level extraction fails.
        text_candidates = re.findall(rb"[A-Za-z0-9][A-Za-z0-9 ,.;:()\\-]{12,}", raw_pdf)
        if text_candidates:
            merged = "\n".join(normalize_text(item.decode("latin-1", errors="ignore")) for item in text_candidates[:40])
            merged = normalize_text(merged)
            if merged:
                units.append(ExtractedUnit(text=merged, metadata={"page_number": 1}))
                warnings.append("pdf_binary_fallback_text")

    return ExtractedDocument(doc_type="pdf", units=units, warnings=warnings)


def extract_pdf(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    try:
        import PyPDF2
    except Exception as exc:  # pragma: no cover
        raise MissingDependencyError("PyPDF2 is required for PDF ingestion") from exc

    raw_pdf = raw_bytes if raw_bytes is not None else read_bytes(path, max_bytes=context.max_bytes)
    reader = PyPDF2.PdfReader(io.BytesIO(raw_pdf))
    return _extract_pdf_from_reader(reader, context, raw_pdf)


def extract_docx(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    try:
        import docx  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise MissingDependencyError("python-docx is required for .docx ingestion") from exc

    raw = _load_raw(path, raw_bytes, context)
    zip_warning = check_zip_safety(
        raw,
        max_entries=context.max_zip_entries,
        max_uncompressed_bytes=context.max_zip_uncompressed_bytes,
    )
    warnings = [zip_warning] if zip_warning else []
    if zip_warning:
        return ExtractedDocument(doc_type="docx", units=[], warnings=warnings)

    with _bytes_as_temp_file(raw_bytes, suffix=".docx") as tmp_path:
        source_path = tmp_path or path
        document = docx.Document(source_path)
        paragraphs = [normalize_text(p.text) for p in document.paragraphs if normalize_text(p.text)]
    units = [ExtractedUnit(text=text) for text in paragraphs]
    return ExtractedDocument(doc_type="docx", units=units, warnings=warnings)


def extract_pptx(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise MissingDependencyError("python-pptx is required for .pptx ingestion") from exc

    raw = _load_raw(path, raw_bytes, context)
    zip_warning = check_zip_safety(
        raw,
        max_entries=context.max_zip_entries,
        max_uncompressed_bytes=context.max_zip_uncompressed_bytes,
    )
    warnings = [zip_warning] if zip_warning else []
    if zip_warning:
        return ExtractedDocument(doc_type="pptx", units=[], warnings=warnings)

    with _bytes_as_temp_file(raw_bytes, suffix=".pptx") as tmp_path:
        source_path = tmp_path or path
        presentation = Presentation(source_path)
        units: List[ExtractedUnit] = []
        for idx, slide in enumerate(presentation.slides, start=1):
            if idx > context.max_slides:
                warnings.append(f"pptx_truncated_at_max_slides={context.max_slides}")
                break
            text_parts = []
            for shape in slide.shapes:
                value = getattr(shape, "text", "")
                if value:
                    text_parts.append(value)
            slide_text = normalize_text("\n".join(text_parts))
            if slide_text:
                units.append(ExtractedUnit(text=slide_text, metadata={"slide_number": idx}))

    return ExtractedDocument(doc_type="pptx", units=units, warnings=warnings)


def extract_xlsx(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    try:
        from openpyxl import load_workbook  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise MissingDependencyError("openpyxl is required for .xlsx ingestion") from exc

    raw = _load_raw(path, raw_bytes, context)
    zip_warning = check_zip_safety(
        raw,
        max_entries=context.max_zip_entries,
        max_uncompressed_bytes=context.max_zip_uncompressed_bytes,
    )
    warnings = [zip_warning] if zip_warning else []
    if zip_warning:
        return ExtractedDocument(doc_type="xlsx", units=[], warnings=warnings)

    with _bytes_as_temp_file(raw_bytes, suffix=".xlsx") as tmp_path:
        source_path = tmp_path or path
        workbook = load_workbook(source_path, read_only=True, data_only=True)
        units: List[ExtractedUnit] = []
        for sheet_idx, sheet_name in enumerate(workbook.sheetnames, start=1):
            if sheet_idx > context.max_sheets:
                warnings.append(f"xlsx_truncated_at_max_sheets={context.max_sheets}")
                break
            sheet = workbook[sheet_name]
            for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_idx > context.max_rows:
                    warnings.append(f"xlsx_sheet_{sheet_name}_truncated_at_max_rows={context.max_rows}")
                    break
                values = ["" if cell is None else str(cell) for cell in row]
                if not any(v.strip() for v in values):
                    continue
                row_text = normalize_text(" | ".join(values))
                units.append(
                    ExtractedUnit(
                        text=row_text,
                        metadata={"sheet_name": sheet_name, "row_number": row_idx},
                    )
                )

    return ExtractedDocument(doc_type="xlsx", units=units, warnings=warnings)


def _ole_extract_strings(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> str:
    try:
        import olefile  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise MissingDependencyError("olefile is required for legacy OLE documents") from exc

    if raw_bytes is not None:
        ole = olefile.OleFileIO(io.BytesIO(raw_bytes))
    else:
        ole = olefile.OleFileIO(path)

    strings: List[str] = []
    try:
        for stream_name in ole.listdir(streams=True, storages=False):
            try:
                stream = ole.openstream(stream_name)
                payload = stream.read()
            except Exception:
                continue
            if not payload:
                continue
            try:
                decoded = payload.decode("utf-16-le")
            except Exception:
                decoded = payload.decode("latin-1", errors="ignore")
            for token in re.findall(r"[\w][\w\s,.;:()\-]{5,}", decoded):
                cleaned = normalize_text(token)
                if len(cleaned) >= 6:
                    strings.append(cleaned)
    finally:
        ole.close()

    return "\n".join(strings)


def extract_doc(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    if not context.enable_legacy_office:
        return ExtractedDocument(doc_type="doc", units=[], warnings=["legacy_office_disabled_by_config"])
    try:
        text = _ole_extract_strings(path, raw_bytes, context)
    except MissingDependencyError:
        raise
    except Exception as exc:
        return ExtractedDocument(doc_type="doc", units=[], warnings=[f"doc_extract_failed: {exc}"])

    if not text.strip():
        return ExtractedDocument(doc_type="doc", units=[], warnings=["doc_no_extractable_text"])
    return ExtractedDocument(doc_type="doc", units=[ExtractedUnit(text=text)], warnings=[])


def extract_ppt(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    if not context.enable_legacy_office:
        return ExtractedDocument(doc_type="ppt", units=[], warnings=["legacy_office_disabled_by_config"])
    try:
        text = _ole_extract_strings(path, raw_bytes, context)
    except MissingDependencyError:
        raise
    except Exception as exc:
        return ExtractedDocument(doc_type="ppt", units=[], warnings=[f"ppt_extract_failed: {exc}"])

    if not text.strip():
        return ExtractedDocument(doc_type="ppt", units=[], warnings=["ppt_no_extractable_text"])
    units = [
        ExtractedUnit(text=chunk, metadata={"slide_number": idx + 1})
        for idx, chunk in enumerate(text.split("\n\n"))
        if chunk.strip()
    ]
    return ExtractedDocument(doc_type="ppt", units=units, warnings=[])


def extract_xls(path: str, raw_bytes: Optional[bytes], context: ExtractorContext) -> ExtractedDocument:
    if not context.enable_legacy_office:
        return ExtractedDocument(doc_type="xls", units=[], warnings=["legacy_office_disabled_by_config"])

    try:
        import xlrd  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise MissingDependencyError("xlrd is required for .xls ingestion") from exc

    try:
        if raw_bytes is not None:
            workbook = xlrd.open_workbook(file_contents=raw_bytes)
        else:
            workbook = xlrd.open_workbook(path)
    except Exception as exc:
        return ExtractedDocument(doc_type="xls", units=[], warnings=[f"xls_open_failed: {exc}"])

    units: List[ExtractedUnit] = []
    warnings: List[str] = []
    for sheet_idx, sheet in enumerate(workbook.sheets(), start=1):
        if sheet_idx > context.max_sheets:
            warnings.append(f"xls_truncated_at_max_sheets={context.max_sheets}")
            break
        for row_idx in range(sheet.nrows):
            if row_idx + 1 > context.max_rows:
                warnings.append(f"xls_sheet_{sheet.name}_truncated_at_max_rows={context.max_rows}")
                break
            values = [str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols)]
            if not any(v.strip() for v in values):
                continue
            units.append(
                ExtractedUnit(
                    text=normalize_text(" | ".join(values)),
                    metadata={"sheet_name": sheet.name, "row_number": row_idx + 1},
                )
            )

    return ExtractedDocument(doc_type="xls", units=units, warnings=warnings)
